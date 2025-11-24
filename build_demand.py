from pathlib import Path
from typing import List

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

GUS_FILE = Path("raporty") / "demografia_dzieci.xlsx"
OUT_XLSX = Path("raporty") / "zapotrzebowanie_miejsc_2023_2060.xlsx"
OUT_PPTX = Path("raporty") / "prezentacja_demografia_placowki.pptx"


def load_powiat():
    df = pd.read_excel(GUS_FILE, sheet_name="powiat_raciborski")
    pivot = df.pivot_table(index="rok", columns="grupa", values="liczba", aggfunc="sum").reset_index()
    pivot = pivot.rename(columns={"zlobek_0_2": "dzieci_0_2", "przedszkole_3_6": "dzieci_3_6", "szkolne_7_18": "dzieci_7_18"})
    for col in ["dzieci_0_2", "dzieci_3_6", "dzieci_7_18"]:
        if col not in pivot.columns:
            pivot[col] = pd.NA
    pivot["dzieci_lacznie"] = pivot[["dzieci_0_2", "dzieci_3_6", "dzieci_7_18"]].sum(axis=1, skipna=True)
    # Założenie: potrzeba miejsc = 100% populacji
    pivot["miejsca_zlobek"] = pivot["dzieci_0_2"]
    pivot["miejsca_przedszkole"] = pivot["dzieci_3_6"]
    pivot["miejsca_szkola"] = pivot["dzieci_7_18"]
    pivot["miejsca_lacznie"] = pivot["dzieci_lacznie"]
    return pivot.sort_values("rok")


def load_miasto():
    df = pd.read_excel(GUS_FILE, sheet_name="miasto_raciborz")
    pivot = df.pivot_table(index="rok", columns="grupa", values="liczba", aggfunc="sum").reset_index()
    # kolumny dostępne: dzieci_0_9 (brak rozbicia 0-2/3-6), mlodziez_10_19, dzieci_0_17, ogolem
    # Uwaga: brak dokładnego podziału 0-2 / 3-6, pozostawiamy NaN w zapotrzebowaniu szczegółowym
    pivot["dzieci_0_2"] = pd.NA
    pivot["dzieci_3_6"] = pd.NA
    pivot["dzieci_7_18_przybl"] = pivot.get("mlodziez_10_19 (przybliżenie grupy szkolnej)", pivot.get("mlodziez_10_19 (przybliżenie grupy szkolnej)", pd.NA))
    pivot["dzieci_0_17"] = pivot.get("dzieci_0_17 (brak rozbicia na 0-2/3-6/7-17)", pd.NA)
    pivot["dzieci_0_9"] = pivot.get("dzieci_0_9 (brak rozbicia 0-2/3-6)", pivot.get("dzieci_0_9 (brak rozbicia 0-2/3-6)", pd.NA))
    pivot["dzieci_lacznie"] = pivot["dzieci_0_17"].fillna(pivot["dzieci_0_9"])
    pivot["miejsca_zlobek"] = pd.NA
    pivot["miejsca_przedszkole"] = pd.NA
    pivot["miejsca_szkola"] = pivot["dzieci_7_18_przybl"]
    pivot["miejsca_lacznie"] = pivot["dzieci_lacznie"]
    return pivot.sort_values("rok")


def save_excel(powiat: pd.DataFrame, miasto: pd.DataFrame):
    with pd.ExcelWriter(OUT_XLSX, engine="openpyxl") as writer:
        powiat.to_excel(writer, sheet_name="powiat_raciborski", index=False)
        miasto.to_excel(writer, sheet_name="miasto_raciborz", index=False)
        # zestawienie
        powiat_assign = powiat.assign(jednostka="Powiat raciborski")
        miasto_assign = miasto.assign(jednostka="Miasto Racibórz")
        combined = pd.concat([powiat_assign, miasto_assign], ignore_index=True)
        combined.to_excel(writer, sheet_name="zestawienie", index=False)


def add_slide_title(prs, title, subtitle=None):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets: List[str]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
    return slide


def add_chart_slide(prs, title, categories, series_dict, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED, pos="B2"):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_dict.items():
        chart_data.add_series(name, values)
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    return slide


def build_ppt(powiat: pd.DataFrame, miasto: pd.DataFrame):
    prs = Presentation()
    add_slide_title(prs, "Demografia i zapotrzebowanie miejsc", "Racibórz i powiat raciborski, prognoza GUS 2023–2060")
    add_bullet_slide(
        prs,
        "Zakres danych",
        [
            "Powiat raciborski: prognoza 2023–2060, grupy 0–2 / 3–6 / 7–18.",
            "Miasto Racibórz: prognoza 2023–2040, dostępne grupy 0–9, 10–19, 0–17 (brak rozbicia 0–2/3–6).",
            "Założenie: zapotrzebowanie na miejsca = 100% liczebności danej grupy wiekowej.",
            "Dane finansowe placówek dostępne osobno w raport_finansowy_2024.xlsx (koszt/ucznia).",
        ],
    )

    sample_years = [2023, 2030, 2040, 2050, 2060]
    powiat_filtered = powiat[powiat["rok"].isin(sample_years)]
    powiat_filtered = powiat_filtered.fillna(0)
    add_chart_slide(
        prs,
        "Powiat raciborski – zapotrzebowanie miejsc (żłobek/przedszkole/szkoła)",
        powiat_filtered["rok"].tolist(),
        {
            "Żłobek 0–2": powiat_filtered["miejsca_zlobek"].tolist(),
            "Przedszkole 3–6": powiat_filtered["miejsca_przedszkole"].tolist(),
            "Szkoła 7–18": powiat_filtered["miejsca_szkola"].tolist(),
        },
    )

    add_chart_slide(
        prs,
        "Powiat raciborski – łącznie dzieci 0–18",
        powiat_filtered["rok"].tolist(),
        {"0–18": powiat_filtered["miejsca_lacznie"].tolist()},
        chart_type=XL_CHART_TYPE.LINE,
    )

    # Miasto Racibórz – tylko dostępne zakresy
    sample_years_m = [y for y in sample_years if y <= miasto["rok"].max()]
    miasto_f = miasto[miasto["rok"].isin(sample_years_m)].fillna(0)
    add_chart_slide(
        prs,
        "Miasto Racibórz – dostępne grupy (0–9, 10–19, 0–17)",
        miasto_f["rok"].tolist(),
        {
            "0–9 (brak podziału 0–2/3–6)": miasto_f["dzieci_0_9"].tolist() if "dzieci_0_9" in miasto_f else [],
            "10–19 (przybliżenie szkolne)": miasto_f["dzieci_7_18_przybl"].tolist(),
            "0–17 łącznie": miasto_f["dzieci_0_17"].tolist(),
        },
    )

    add_bullet_slide(
        prs,
        "Kluczowe uwagi",
        [
            "Miasto Racibórz: brak rozbicia na 0–2 i 3–6 – zalecane uzupełnienie danymi lokalnymi.",
            "Powiat: pełne rozbicie 0–2/3–6/7–18 dostępne z prognozy GUS 2023–2060.",
            "Kolejne kroki: zestawić z pojemnością placówek (żłobki, przedszkola, szkoły) i kosztami/ucznia.",
        ],
    )
    prs.save(OUT_PPTX)


def main():
    powiat = load_powiat()
    miasto = load_miasto()
    save_excel(powiat, miasto)
    build_ppt(powiat, miasto)
    print(f"Zapisano {OUT_XLSX}")
    print(f"Zapisano {OUT_PPTX}")


if __name__ == "__main__":
    main()
